# -*- coding: utf-8 -*-
"""
elite-powerpoint-designer 风格实现
Corporate Professional (Microsoft/IBM Style) 品牌系统
藏青 Navy #003366 / 钢蓝 Steel Blue #0078D4 / 暖灰 Warm Gray #F3F2F1
设计原则：极简、大胆清晰、视觉层级、一致性、有目的的动效
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- 品牌系统 ----------
NAVY      = RGBColor(0x00, 0x33, 0x66)   # primary
STEEL     = RGBColor(0x00, 0x78, 0xD4)   # accent
WARM_GRAY = RGBColor(0xF3, 0xF2, 0xF1)   # bg tint
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x22, 0x22, 0x22)   # body text
MID       = RGBColor(0x59, 0x59, 0x59)   # secondary text
LIGHT     = RGBColor(0xA6, 0xA6, 0xA6)   # caption
NAVY_DEEP = RGBColor(0x00, 0x24, 0x4C)   # title bg darker

FONT = "微软雅黑"

# 16:9
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.65)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

# ---------- 工具函数 ----------

def set_run_font(run, size=Pt(16), bold=False, color=DARK, name=FONT):
    f = run.font
    f.name = name
    f.size = size
    f.bold = bold
    f.color.rgb = color
    # 同时设置东亚字体，保证中文用微软雅黑
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def add_text(slide, x, y, w, h, text, size=Pt(16), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
             space_after=Pt(0)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = space_after
        run = p.add_run()
        run.text = ln
        set_run_font(run, size, bold, color)
    return tb

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75),
             shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.shadow.inherit = False
    return sp

def add_shape_text(sp, text, size=Pt(14), bold=False, color=DARK,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        set_run_font(run, size, bold, color)
    return sp

def add_slide_blank():
    return prs.slides.add_slide(BLANK)

def bg_fill(slide, color):
    add_rect(slide, 0, 0, SW, SH, fill=color)

def add_page_header(slide, kicker, title, title_color=DARK, accent=STEEL):
    """kicker: 小节标签(小号藏青)；title: 页面主标题"""
    add_text(slide, MARGIN, Inches(0.42), Inches(9.5), Inches(0.34),
             kicker, size=Pt(12), bold=True, color=accent)
    add_text(slide, MARGIN, Inches(0.78), Inches(12), Inches(0.8),
             title, size=Pt(30), bold=True, color=title_color)
    # 标题下细基线（数据友好型布局）
    add_rect(slide, MARGIN, Inches(1.62), Inches(1.1), Pt(3.2), fill=accent)
    add_rect(slide, MARGIN + Inches(1.1), Inches(1.62), Inches(4.6), Pt(1.0), fill=LIGHT)

def metric_card(slide, x, y, w, h, value, label, accent=STEEL, dark=False):
    """KPI 卡片：大数字 + 小标签"""
    card = add_rect(slide, x, y, w, h, fill=(NAVY if dark else WARM_GRAY),
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_shape_text(card, value, size=Pt(26), bold=True, color=(WHITE if dark else accent))
    # 标签放在卡片底部，作为第二段
    tf = card.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    set_run_font(run, size=Pt(10.5), bold=False, color=(WHITE if dark else MID))
    return card

def fact_item(slide, x, y, w, text, size=Pt(13.5), color=DARK, marker="✓", marker_color=STEEL):
    """事实条目：左侧钢蓝对勾 + 文本"""
    add_text(slide, x, y + Inches(0.03), Inches(0.4), Inches(0.4),
             marker, size=Pt(15), bold=True, color=marker_color)
    tb = add_text(slide, x + Inches(0.42), y, w, Inches(0.4), text,
                  size=size, color=color, line_spacing=1.18)
    # 返回高度估算用的文本对象
    return tb

# ============================================================
# Slide 1 · 封面（title_slide / hero）
# ============================================================
s = add_slide_blank()
bg_fill(s, NAVY_DEEP)
add_rect(s, 0, 0, SW, Inches(0.10), fill=STEEL)
# 顶部小标签
add_text(s, MARGIN, Inches(0.85), Inches(10), Inches(0.4),
         "MELONAI · 职场AI日报 · 2026-08-02", size=Pt(13), bold=True, color=STEEL)
# 主标题
add_text(s, MARGIN, Inches(1.55), Inches(12.0), Inches(2.2),
         "今日三大AI要闻", size=Pt(54), bold=True, color=WHITE, line_spacing=1.05)
add_text(s, MARGIN, Inches(3.35), Inches(12.0), Inches(0.9),
         "打工人视角解读 —— 5分钟看懂今天该知道的 AI 大事", size=Pt(22), color=RGBColor(0xBF,0xD4,0xEA))
# 三件事预览（底部三列）
labels = [
    ("01", "国产免费反超", "DeepSeek-V4-Flash 正式版，干活能力接近顶级"),
    ("02", "AI 干长活立功", "Astra 花 $2000 连破 10 道数学世纪难题"),
    ("03", "开源登顶·资本大战", "Kimi K3 全球最大开源模型，背后估值350亿"),
]
col_w = Inches(3.85)
gap = (SW - 2*MARGIN - 3*col_w) / 2
for i, (num, t, d) in enumerate(labels):
    x = MARGIN + i*(col_w + gap)
    y = Inches(4.75)
    card = add_rect(s, x, y, col_w, Inches(1.85), fill=RGBColor(0x0A,0x41,0x73),
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    add_shape_text(card, "", size=Pt(1))
    add_text(s, x + Inches(0.28), y + Inches(0.22), Inches(0.8), Inches(0.5),
             num, size=Pt(26), bold=True, color=STEEL)
    add_text(s, x + Inches(0.28), y + Inches(0.85), col_w - Inches(0.5), Inches(0.45),
             t, size=Pt(15), bold=True, color=WHITE)
    add_text(s, x + Inches(0.28), y + Inches(1.28), col_w - Inches(0.5), Inches(0.5),
             d, size=Pt(10.5), color=RGBColor(0xBF,0xD4,0xEA))
add_text(s, MARGIN, Inches(7.05), Inches(10), Inches(0.35),
         "来源：官方博客 / 官方 GitHub / 官方产品站 ｜ 融资信息均标注「据媒体报道」",
         size=Pt(10), color=RGBColor(0x7A,0x9A,0xC0))

# ============================================================
# Slide 2 · 目录（agenda / 三卡片）
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "AGENDA · 今日目录", "今天要聊的三件大事")
agenda = [
    ("01", "DeepSeek-V4-Flash 正式版发布", "国产免费模型反超：顶级干活能力，几乎免费",
     "总参 284B / 激活 13B · 1M 上下文 · 单Token算力前代27%"),
    ("02", "OpenAI Astra 数学突破", "AI 花 $2000 连破 10 道数学世纪难题",
     "249 页手稿 · Lean 4 形式化验证 · next major model"),
    ("03", "Kimi K3 开源 + 融资", "全球最大开源模型免费给你用，背后公司要上市",
     "2.8T 总参 · F轮35亿美元 · 估值350亿 · 港股传闻"),
]
for i, (num, t, d, sub) in enumerate(agenda):
    y = Inches(2.05) + i * Inches(1.72)
    card = add_rect(s, MARGIN, y, Inches(12.03), Inches(1.5),
                    fill=WARM_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    # 左侧数字块
    add_rect(s, MARGIN, y, Inches(1.5), Inches(1.5), fill=NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_text(s, MARGIN, y + Inches(0.45), Inches(1.5), Inches(0.6),
             num, size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN + Inches(1.85), y + Inches(0.20), Inches(9.9), Inches(0.5),
             t, size=Pt(17), bold=True, color=NAVY)
    add_text(s, MARGIN + Inches(1.85), y + Inches(0.72), Inches(9.9), Inches(0.4),
             d, size=Pt(12.5), color=DARK)
    add_text(s, MARGIN + Inches(1.85), y + Inches(1.10), Inches(9.9), Inches(0.35),
             sub, size=Pt(10.5), color=MID)

# ============================================================
# Slide 3 · 话题一 核心事实（metrics + facts）
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "TOPIC 01 · DEEPSEEK-V4-FLASH", "国产免费模型，追到顶级 AI 脸上了")
# 左侧事实
facts = [
    "2026-07-31 正式版 API 上线公测，现有集成不改代码自动升级",
    "轻量化 MoE：总参数 284B / 激活 13B，1M 超长上下文",
    "官方 9 项 Agent 基准多项超过三个月前的 V4-Pro 预览版",
    "原生兼容 OpenAI Responses API，针对性适配 Codex",
    "支持思考 / 非思考模式切换",
    "百万 Token 场景单 Token 算力仅为前代 27%，价格极低",
]
y = Inches(2.0)
for i, f in enumerate(facts):
    fact_item(s, MARGIN, y + i*Inches(0.62), Inches(6.6), f, size=Pt(12.5))
# 右侧 KPI 面板
panel_x = Inches(7.55)
panel = add_rect(s, panel_x, Inches(1.95), Inches(5.15), Inches(4.75),
                 fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_text(s, panel_x + Inches(0.3), Inches(2.2), Inches(4.5), Inches(0.4),
         "核心数据", size=Pt(15), bold=True, color=WHITE)
metrics = [
    ("284B", "总参数"), ("13B", "激活参数"), ("1M", "上下文长度"),
    ("27%", "单Token算力成本"),
    ("$0.0028", "缓存命中输入价(每百万token)"),
]
for i, (v, l) in enumerate(metrics):
    col = i % 2
    row = i // 2
    mx = panel_x + Inches(0.3) + col * Inches(2.42)
    my = Inches(2.72) + row * Inches(1.18)
    card = add_rect(s, mx, my, Inches(2.25), Inches(1.0),
                    fill=RGBColor(0x0A,0x41,0x73), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    add_text(s, mx, my + Inches(0.12), Inches(2.25), Inches(0.5),
             v, size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, mx, my + Inches(0.60), Inches(2.25), Inches(0.35),
             l, size=Pt(10), color=RGBColor(0xBF,0xD4,0xEA), align=PP_ALIGN.CENTER)

# ============================================================
# Slide 4 · 话题一 打工人视角 + 行动
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "TOPIC 01 · 打工人视角", "跟你有什么关系？")
# 视角卡
vx, vy, vw, vh = MARGIN, Inches(1.95), Inches(12.03), Inches(2.6)
view = add_rect(s, vx, vy, vw, vh, fill=WARM_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
add_rect(s, vx, vy, Inches(0.09), vh, fill=STEEL)
add_text(s, vx + Inches(0.45), vy + Inches(0.3), Inches(11.2), Inches(0.45),
         "一句话翻译", size=Pt(15), bold=True, color=STEEL)
add_text(s, vx + Inches(0.45), vy + Inches(0.85), Inches(11.2), Inches(1.6),
         "顶级干活能力 + 极低成本 → 普通打工人也能用上顶尖模型。\n"
         "选工具不用迷信国外付费产品；拆掉「必须花钱买课 / 买会员才能用好 AI」的焦虑。",
         size=Pt(15.5), color=DARK, line_spacing=1.35)
# 行动卡
ax, ay, aw, ah = MARGIN, Inches(4.9), Inches(12.03), Inches(1.95)
act = add_rect(s, ax, ay, aw, ah, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
add_text(s, ax + Inches(0.45), ay + Inches(0.28), Inches(11.2), Inches(0.45),
         "今天就能带走 · 行动", size=Pt(15), bold=True, color=RGBColor(0x9E,0xC6,0xEE))
add_text(s, ax + Inches(0.45), ay + Inches(0.8), Inches(11.2), Inches(1.0),
         "打开 chat.deepseek.com 或 DeepSeek 开放平台，用 V4-Flash 干一件真实工作上的活\n"
         "（写周报 / 整理数据 / 写代码），和现在用的工具对比一次。",
         size=Pt(14.5), color=WHITE, line_spacing=1.3)

# ============================================================
# Slide 5 · 话题二 核心事实
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "TOPIC 02 · OPENAI ASTRA", "AI 花 $2000，连破 10 道数学世纪难题")
facts = [
    "OpenAI 官方博客首次公开下一代模型 Astra（next major model）内部版本",
    "在 10 个长期未决的数学 / 理论计算机科学难题上做出新结果",
    "领域：高维几何、编码理论、群论、量子复杂性、格密码学等",
    "每个论证用 Lean 4 形式化验证，证书公开在 GitHub",
    "发布 249 页手稿合集 + 每个问题的 AI 推理过程旁白",
    "官方称全部 token 成本约 2000 美元",
]
y = Inches(2.0)
for i, f in enumerate(facts):
    fact_item(s, MARGIN, y + i*Inches(0.62), Inches(6.6), f, size=Pt(12.5))
# 右侧 KPI
panel_x = Inches(7.55)
panel = add_rect(s, panel_x, Inches(1.95), Inches(5.15), Inches(4.75),
                 fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_text(s, panel_x + Inches(0.3), Inches(2.2), Inches(4.5), Inches(0.4),
         "核心数据", size=Pt(15), bold=True, color=WHITE)
metrics = [
    ("10", "解出的世纪难题"), ("249页", "公开手稿"), ("$2000", "Token 成本"),
    ("Lean4", "形式化验证"), ("0", "人类卡题年限下限"),
]
for i, (v, l) in enumerate(metrics):
    col = i % 2
    row = i // 2
    mx = panel_x + Inches(0.3) + col * Inches(2.42)
    my = Inches(2.72) + row * Inches(1.18)
    card = add_rect(s, mx, my, Inches(2.25), Inches(1.0),
                    fill=RGBColor(0x0A,0x41,0x73), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    add_text(s, mx, my + Inches(0.12), Inches(2.25), Inches(0.5),
             v, size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, mx, my + Inches(0.60), Inches(2.25), Inches(0.35),
             l, size=Pt(10), color=RGBColor(0xBF,0xD4,0xEA), align=PP_ALIGN.CENTER)

# ============================================================
# Slide 6 · 话题二 打工人视角 + 行动
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "TOPIC 02 · 打工人视角", "AI 从「答一句」升级到「干一件长活」")
vx, vy, vw, vh = MARGIN, Inches(1.95), Inches(12.03), Inches(2.6)
view = add_rect(s, vx, vy, vw, vh, fill=WARM_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
add_rect(s, vx, vy, Inches(0.09), vh, fill=STEEL)
add_text(s, vx + Inches(0.45), vy + Inches(0.3), Inches(11.2), Inches(0.45),
         "一句话翻译", size=Pt(15), bold=True, color=STEEL)
add_text(s, vx + Inches(0.45), vy + Inches(0.85), Inches(11.2), Inches(1.6),
         "数学难题本身离打工人远，但内核是：AI 从「答一句」升级到「独立干一件长活」。\n"
         "以后你可以把整件工作交给 AI，而不是一句一句喂。\n"
         "接续上一期「AI 自己干活没人盯会闯祸」——这次是「AI 干长活立了功」。",
         size=Pt(15.5), color=DARK, line_spacing=1.35)
ax, ay, aw, ah = MARGIN, Inches(4.9), Inches(12.03), Inches(1.95)
act = add_rect(s, ax, ay, aw, ah, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
add_text(s, ax + Inches(0.45), ay + Inches(0.28), Inches(11.2), Inches(0.45),
         "今天就能带走 · 行动", size=Pt(15), bold=True, color=RGBColor(0x9E,0xC6,0xEE))
add_text(s, ax + Inches(0.45), ay + Inches(0.8), Inches(11.2), Inches(1.0),
         "把一个能拆成「交代清楚就能跑」的小任务（如让它自己调研 + 整理成报告），\n"
         "整件委托给 AI Agent 试一次，观察它能自主跑多久、做到什么程度。",
         size=Pt(14.5), color=WHITE, line_spacing=1.3)

# ============================================================
# Slide 7 · 话题三 核心事实
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "TOPIC 03 · KIMI K3 开源 + 融资", "全球最大开源 AI 模型，免费给你用")
facts = [
    "月之暗面发布 Kimi K3：全球首个开源的三万亿级别大模型",
    "7/27 全链条开源：权重 + 47 页技术报告 + 三套基础设施",
    "WebDev Arena 以 1678 Elo 登顶，首个登顶的开源模型",
    "7/29 完成超 35 亿美元 F 轮，投后估值 350 亿美元",
    "F 轮认购超 3 倍提前关闭，市场传闻最快 6 个月港股上市",
    "1M token 上下文，原生视觉 MoonViT-V2",
]
y = Inches(2.0)
for i, f in enumerate(facts):
    fact_item(s, MARGIN, y + i*Inches(0.62), Inches(6.6), f, size=Pt(12.5))
panel_x = Inches(7.55)
panel = add_rect(s, panel_x, Inches(1.95), Inches(5.15), Inches(4.75),
                 fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
add_text(s, panel_x + Inches(0.3), Inches(2.2), Inches(4.5), Inches(0.4),
         "核心数据", size=Pt(15), bold=True, color=WHITE)
metrics = [
    ("2.8T", "总参数(全球首个3万亿级)"), ("104B", "激活参数"),
    ("1M", "上下文长度"), ("$35亿", "F轮融资"),
    ("$350亿", "投后估值"),
]
for i, (v, l) in enumerate(metrics):
    col = i % 2
    row = i // 2
    mx = panel_x + Inches(0.3) + col * Inches(2.42)
    my = Inches(2.72) + row * Inches(1.18)
    card = add_rect(s, mx, my, Inches(2.25), Inches(1.0),
                    fill=RGBColor(0x0A,0x41,0x73), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    add_text(s, mx, my + Inches(0.12), Inches(2.25), Inches(0.5),
             v, size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, mx, my + Inches(0.60), Inches(2.25), Inches(0.35),
             l, size=Pt(10), color=RGBColor(0xBF,0xD4,0xEA), align=PP_ALIGN.CENTER)

# ============================================================
# Slide 8 · 话题三 打工人视角 + 行动
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "TOPIC 03 · 打工人视角", "看懂资本，才不会被「焦虑营销」收割")
vx, vy, vw, vh = MARGIN, Inches(1.95), Inches(12.03), Inches(2.6)
view = add_rect(s, vx, vy, vw, vh, fill=WARM_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
add_rect(s, vx, vy, Inches(0.09), vh, fill=STEEL)
add_text(s, vx + Inches(0.45), vy + Inches(0.3), Inches(11.2), Inches(0.45),
         "一句话翻译", size=Pt(15), bold=True, color=STEEL)
add_text(s, vx + Inches(0.45), vy + Inches(0.85), Inches(11.2), Inches(1.6),
         "① 顶级能力免费开源 → AI 只会越来越便宜，49 块的 AI 课真没必要买；\n"
         "② 你天天用的 Kimi，背后是几百亿美元资本大战——\n"
         "　免费背后有人替你付费，看懂资本才不会被「焦虑营销」收割。",
         size=Pt(15.5), color=DARK, line_spacing=1.35)
ax, ay, aw, ah = MARGIN, Inches(4.9), Inches(12.03), Inches(1.95)
act = add_rect(s, ax, ay, aw, ah, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
add_text(s, ax + Inches(0.45), ay + Inches(0.28), Inches(11.2), Inches(0.45),
         "今天就能带走 · 行动", size=Pt(15), bold=True, color=RGBColor(0x9E,0xC6,0xEE))
add_text(s, ax + Inches(0.45), ay + Inches(0.8), Inches(11.2), Inches(1.0),
         "打开 Kimi（kimi.com）用 K3 试一个你正犹豫要不要付费的功能，\n"
         "对比一下免费开源到底能不能打。",
         size=Pt(14.5), color=WHITE, line_spacing=1.3)

# ============================================================
# Slide 9 · 今日行动清单（checklist）
# ============================================================
s = add_slide_blank()
bg_fill(s, WHITE)
add_page_header(s, "SUMMARY · 三句话带走", "今天下班就能做")
takeaways = [
    ("工具认知更新", "国产免费模型已具备顶级干活能力，选 AI 工具不用迷信国外付费"),
    ("使用方式升级", "AI 能独立干一整件长活了：从「一句一句喂」变成「整件委托」"),
    ("消费决策清醒", "AI 只会越来越便宜，别为焦虑买单，看懂免费背后的资本逻辑"),
]
for i, (t, d) in enumerate(takeaways):
    y = Inches(2.05) + i * Inches(1.15)
    card = add_rect(s, MARGIN, y, Inches(12.03), Inches(0.98),
                    fill=WARM_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    add_rect(s, MARGIN, y, Inches(0.09), Inches(0.98), fill=STEEL)
    add_text(s, MARGIN + Inches(0.4), y + Inches(0.16), Inches(2.7), Inches(0.45),
             "▢ " + t, size=Pt(14), bold=True, color=NAVY)
    add_text(s, MARGIN + Inches(3.2), y + Inches(0.16), Inches(8.6), Inches(0.7),
             d, size=Pt(12.5), color=DARK)
actions = [
    "用 DeepSeek V4-Flash 干一件真实工作上的活，和现用工具对比",
    "把一个「交代清楚就能跑」的小任务整件委托给 AI Agent",
    "打开 Kimi 用 K3 试一个犹豫要不要付费的功能",
]
ay = Inches(5.55)
add_text(s, MARGIN, ay, Inches(12), Inches(0.4),
         "三件小事，今天下班就能做", size=Pt(15), bold=True, color=STEEL)
for i, a in enumerate(actions):
    fact_item(s, MARGIN, ay + Inches(0.5) + i*Inches(0.42), Inches(11.5),
              a, size=Pt(12.5), marker="▢")

# ============================================================
# Slide 10 · 结尾（closing / navy）
# ============================================================
s = add_slide_blank()
bg_fill(s, NAVY_DEEP)
add_rect(s, 0, 0, SW, Inches(0.10), fill=STEEL)
add_text(s, MARGIN, Inches(2.35), Inches(12), Inches(1.1),
         "职场AI · 每天 5 分钟，读懂 AI 大事", size=Pt(40), bold=True, color=WHITE)
add_text(s, MARGIN, Inches(3.7), Inches(12), Inches(0.6),
         "下期见", size=Pt(20), color=RGBColor(0xBF,0xD4,0xEA))
add_text(s, MARGIN, Inches(6.4), Inches(12), Inches(0.4),
         "数据来源：官方博客 / 官方 GitHub / 官方产品站 ｜ 融资类信息均标注「据媒体报道」",
         size=Pt(10.5), color=RGBColor(0x7A,0x9A,0xC0))

# ---------- 保存 ----------
out = r"d:/dataspath/Claude/video-script-generator/PPT生成/1-elite-powerpoint-designer/今日三大AI要闻_elite.pptx"
prs.save(out)
print("saved:", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
